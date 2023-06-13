""""The main app of my flask proyect"""
import json
import os
import openai
import requests
import urllib.parse
from flask import Flask, make_response, render_template, request, send_file,redirect,jsonify
from pptx.util import Inches,Pt
from pptx import Presentation

app = Flask('run',template_folder='app/views')   

def chatGPTrequest(text):
    openai.api_key  = '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'
    prompt = 'Make me a JSON file of a PowerPoint presentation of the theme "'+text+'" in Spanish. This JSON file will contain multiple slides, each slide possibly containing the following: Title of the slide, search querys of images it may have, and possible text.Here is an example about deer:{"slides": [{"title": "Definición de venado","image": "deer","text": "Los venados son animales majestuosos que habitan en diversas regiones del mundo."},{"title": "Características del venado","image": "different species of deer","text": "Existen diferentes especies de venados, cada una con características únicas."},{"title": "Cuernos ramificados del venado","image": "deer horns","text": "Los venados poseen cuernos ramificados que son utilizados para la competencia y atracción de parejas."}, {"title": "Hábitats del venado","image": "A deer in a habitat","text": "Los venados se adaptan a diferentes hábitats, desde bosques hasta praderas y montañas."},{"title": "Importancia de la conservación del venado","image": "The importance of deer","text": "La conservación de los venados es crucial para preservar la biodiversidad y el equilibrio de los ecosistemas."}]}And so on, with multiple slides. Please note that:1. I want the image attribute to contain a simple search query in english of an image for pixabay not a URL.2.The image search query should not exceed 100 characters.3. The title attribute should contain the title of the slide, not the number of the slide. 4.I only want the JSON file in your response. Do not include any additional text, only JSON.'
    response = openai.Completion.create(
        engine='text-davinci-003',
        prompt=prompt,
        max_tokens=1000,
        n=1,
        stop=None,
        temperature=0.7
    )
    if 'choices' in response and len(response['choices']) > 0:
        return response['choices'][0]['text']
    return {}

@app.route('/')
def index(): 
     return render_template('index.html')

@app.route('/download')
def descargar():
    json_data = chatGPTrequest(request.args.get('input_text'))
    print(f"PROMPT\n{json_data}")
    presentation = create_powerpoint(json_data)
    presentation_path = 'output.pptx'
    presentation.save(presentation_path)

    # Prepare the file to be downloaded
    with open(presentation_path, 'rb') as file:
        data = file.read()

    response = make_response(data)
    response.headers['Content-Disposition'] = 'attachment; filename=output.pptx'
    response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

    return response

def search_image(query):
    api_key = '37240249-b6a284da7936917de819e23e0'
    url = f"https://pixabay.com/api/?key={api_key}&q={urllib.parse.quote(query)}"
    response = requests.get(url,timeout=5)
    data = response.json()
    if 'hits' in data:
        if len(data['hits']) > 0:
            image_url = data['hits'][0]['largeImageURL']
            return image_url

    return None
def download_image(url):
    response = requests.get(url,timeout=5)
    if response.status_code == 200:
        image_data = response.content
        image_path = 'temp.jpg'
        with open(image_path, 'wb') as file:
            file.write(image_data)
        return image_path
    return None


def create_powerpoint(json_data):
    presentation = Presentation()

    data = json.loads(json_data)
    if 'slides' not in data:
        # Handle case where 'slides' key is missing
        return None

    for slide_data in data['slides']:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        # Set slide title
        title = slide.shapes.title
        title.text = slide_data.get('title', '')

        # Set slide text
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        text_frame = content_box.text_frame
        text_frame.text = slide_data.get('text', '')

        # Configure text formatting
        text_frame.word_wrap = True
        text_frame.auto_size = True

        # Search image on Pixabay
        search_query = slide_data.get('image', '')
        image_url = search_image(search_query)

        # Download image and add to slide
        if image_url:
            image_path = download_image(image_url)
            if image_path:
                image_box = slide.shapes.add_picture(image_path, Inches(1), Inches(3), width=Inches(6), height=Inches(4))

                # Adjust image positioning
                image_box.left = int((presentation.slide_width - image_box.width) / 2)

            # Delete downloaded image from server
            if os.path.exists(image_path):
                os.remove(image_path)

        # Adjust font size for longer text
        if len(text_frame.text) > 200:
            text_frame.text = slide_data.get('text', '')
            font = text_frame.paragraphs[0].runs[0].font
            font.size = Pt(16)

    return presentation
if __name__ == '__main__':
    app.run(debug=True)

