
import os
import json
import random
from pptx.util import Inches,Pt
from pptx import Presentation
from app.controller.imagesController import imageController

              #Template img Dimensions:
              #Atlas : 7.2, 4.4 w = 3, h = 2 
              #Estela: 4.1 3.75 w = 5 h = 3.33
              #Citables 4.96 4.92  w = 3.5  h = 2.33
class powerPointController:
    def __init__(self):
        self.imgController = imageController()

    def chooseTemplate(self):
        templates = ["Atlas", "Estela","Citables"]
        usedTemplate = random.choice(templates)

        info = {
            "Atlas": {"name":"Atlas","vertical": 7.2, "horizontal":4.4, "width":3,"height":2},
            "Estela":{"name":"Estela","vertical": 4.17, "horizontal":3.75, "width":5,"height":3.33},
            "Citables":{"name":"Citables","vertical": 4.96, "horizontal":4.92, "width":3.5,"height":2.33}
        }
        return info[usedTemplate]



    def createPowerpoint(self,json_data):
        template = self.chooseTemplate()
        name = template["name"]
        vertical = template["vertical"]
        horizontal = template["horizontal"]
        width = template["width"]
        height = template["height"]
        presentation = Presentation(f"app/templates/{name}.pptx")

        data = json.loads(json_data)
        if 'slides' not in data:
            # Handle case where 'slides' key is missing
            return None

        for slide_data in data['slides']:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])

            # Set slide title
            title = slide.shapes[0]
            title.text = slide_data.get('title', '')

            # Set slide text
            text_frame = slide.shapes[1]
            text_frame.text = slide_data.get('text', '')

            # Configure text formatting
            text_frame.word_wrap = True
            text_frame.auto_size = True

            # Search image on Pixabay
            search_query = slide_data.get('image', '')
            image_url = self.imgController.searchImage(search_query)

            # Download image and add to slide
            if image_url:
                image_path = self.imgController.downloadImage(image_url)
                if image_path:
                    slide.shapes.add_picture(image_path, Inches(vertical), Inches(horizontal), width=Inches(width), height=Inches(height))


                # Delete downloaded image from server
                if os.path.exists(image_path):
                    os.remove(image_path)

            # Adjust font size for longer text
            if len(text_frame.text) > 200:
                text_frame.text = slide_data.get('text', '')
                font = text_frame.paragraphs[0].runs[0].font
                font.size = Pt(16)

        return presentation